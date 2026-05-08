"""
Insert "Strategy Validation & Risk Disclosure" slide at position 11
(after slide 10 Walk-Forward Validation) into HMM_Strategy_Pitch.pptx
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PPTX_PATH = "/Users/zhaowenxuan/Desktop/工作/HMM/HMM_Strategy_Pitch.pptx"

# ─── Color constants ──────────────────────────────────────────────────────────
BG_COLOR    = RGBColor(0x0F, 0x11, 0x17)
GOLD        = RGBColor(0xFF, 0xD7, 0x40)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB4, 0xB4, 0xB4)
DARK_GOLD   = RGBColor(0x50, 0x41, 0x0A)
DIM_GRAY    = RGBColor(0x64, 0x64, 0x64)
GREEN_LIGHT = RGBColor(0xA8, 0xE6, 0xA3)
AMBER_LIGHT = RGBColor(0xFF, 0xD0, 0x8A)

# ─── Slide dimensions (from source file) ─────────────────────────────────────
SW = 12191695   # slide width  (EMU)
SH = 6858000    # slide height (EMU)


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def add_solid_fill(shape, color: RGBColor):
    """Apply solid fill to a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height):
    """Add a transparent text box and return it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.background()
    return txBox


def set_run(run, text, font_size_pt, bold=False, color=WHITE, italic=False):
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_paragraph(tf, text, font_size_pt, bold=False, color=WHITE,
                  italic=False, space_before_pt=0, space_after_pt=0,
                  align=PP_ALIGN.LEFT):
    """Add a paragraph with a single run to a text frame."""
    para = tf.add_paragraph()
    para.alignment = align
    pPr = para._pPr
    if pPr is None:
        pPr = para._p.get_or_add_pPr()
    pPr.set(qn("a:spcBef"), "")
    # Use lxml to set spacing
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", str(int(space_before_pt * 100)))

    run = para.add_run()
    set_run(run, text, font_size_pt, bold=bold, color=color, italic=italic)
    return para


# ─── Main ─────────────────────────────────────────────────────────────────────
def build_slide(prs: Presentation):
    """Build the new slide and return it (appended at end for now)."""
    blank_layout = prs.slide_layouts[6]   # "Blank" layout
    slide = prs.slides.add_slide(blank_layout)

    # ── Background ────────────────────────────────────────────────────────────
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # ── Top gold bar (full width, ~4pt tall = ~5333 EMU per pt → 4pt ≈ 50800) ─
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left=0, top=0, width=SW, height=Pt(4)
    )
    add_solid_fill(bar, GOLD)
    bar.line.fill.background()  # no border

    # ── Title text box ────────────────────────────────────────────────────────
    # English title + Chinese subtitle, mirroring slide 10 style
    title_box = add_text_box(slide,
                             left=Emu(548640), top=Emu(457200),
                             width=Emu(11064240), height=Emu(914400))
    tf = title_box.text_frame
    tf.word_wrap = False

    # First paragraph: English title
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    set_run(r0, "Strategy Validation & Risk Disclosure",
            font_size_pt=30, bold=True, color=GOLD)

    # Second paragraph: Chinese subtitle
    add_paragraph(tf, "策略验证与风险说明",
                  font_size_pt=16, bold=False, color=LIGHT_GRAY,
                  space_before_pt=2)

    # ── Separator line (matches slide 10 Rectangle 3) ─────────────────────────
    sep = slide.shapes.add_shape(
        1,
        left=Emu(548640), top=Emu(1508760),
        width=Emu(11064240), height=Emu(12700)
    )
    add_solid_fill(sep, DARK_GOLD)
    sep.line.fill.background()

    # ────────────────────────────────────────────────────────────────────────
    # Layout parameters
    # Two columns: left starts at ~549k EMU, width ~5100k; right starts ~6250k
    # ────────────────────────────────────────────────────────────────────────
    COL_TOP    = Emu(1600000)
    COL_HEIGHT = Emu(3350000)

    L_LEFT  = Emu(548640)
    L_WIDTH = Emu(5200000)

    R_LEFT  = Emu(6380000)
    R_WIDTH = Emu(5230000)

    BLOCK_GAP = Emu(120000)   # gap between blocks

    # ── LEFT COLUMN ───────────────────────────────────────────────────────────
    left_box = add_text_box(slide,
                            left=L_LEFT, top=COL_TOP,
                            width=L_WIDTH, height=COL_HEIGHT)
    lt = left_box.text_frame
    lt.word_wrap = True

    # Column header
    p = lt.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, "✅  已解决的三个关键问题（Issues Resolved）",
            font_size_pt=13, bold=True, color=GOLD)

    # ── Block 1 ──
    add_paragraph(lt, "",  font_size_pt=6, color=WHITE, space_before_pt=0)   # spacer
    add_paragraph(lt, "◆  模型稳定性 / Model Stability",
                  font_size_pt=11.5, bold=True, color=GREEN_LIGHT, space_before_pt=4)
    add_paragraph(lt,
                  "早期每次训练结果差异达4倍。改为5种随机种子取最优 log-likelihood 模型，结果现已完全可复现。",
                  font_size_pt=10, color=WHITE, space_before_pt=2)
    add_paragraph(lt,
                  "Previously results varied 4× across runs. Now using best of 5 seeds by log-likelihood — fully reproducible.",
                  font_size_pt=9.5, italic=True, color=LIGHT_GRAY, space_before_pt=1)

    # ── Block 2 ──
    add_paragraph(lt, "",  font_size_pt=4, color=WHITE)   # spacer
    add_paragraph(lt, "◆  样本量 / Sample Size",
                  font_size_pt=11.5, bold=True, color=GREEN_LIGHT, space_before_pt=4)
    add_paragraph(lt,
                  "将历史从2016年拉长至2007年，覆盖2008金融危机、2011黄金牛市、2020疫情冲击。Gold交易笔数从16笔增至25笔，AAPL从38笔增至51笔。",
                  font_size_pt=10, color=WHITE, space_before_pt=2)
    add_paragraph(lt,
                  "Extended history from 2016 to 2007, covering 2008 GFC, 2011 Gold bull, 2020 COVID crash. Gold trades: 16→25, AAPL: 38→51.",
                  font_size_pt=9.5, italic=True, color=LIGHT_GRAY, space_before_pt=1)

    # ── Block 3 ──
    add_paragraph(lt, "",  font_size_pt=4, color=WHITE)   # spacer
    add_paragraph(lt, "◆  参数验证 / Parameter Validation",
                  font_size_pt=11.5, bold=True, color=GREEN_LIGHT, space_before_pt=4)
    add_paragraph(lt,
                  "三段子样本验证（2007-2012 / 2013-2018 / 2019-2026），AAPL与Gold三段全部正收益，参数无过拟合。Silver早期两段表现弱，已如实说明。",
                  font_size_pt=10, color=WHITE, space_before_pt=2)
    add_paragraph(lt,
                  "3-period subsample test: AAPL & Gold positive in all 3 periods — no overfitting. Silver underperformed in early periods, disclosed transparently.",
                  font_size_pt=9.5, italic=True, color=LIGHT_GRAY, space_before_pt=1)

    # ── RIGHT COLUMN ──────────────────────────────────────────────────────────
    right_box = add_text_box(slide,
                             left=R_LEFT, top=COL_TOP,
                             width=R_WIDTH, height=COL_HEIGHT)
    rt = right_box.text_frame
    rt.word_wrap = True

    # Column header
    p = rt.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, "⚠️  待改进项（Pending Improvement）",
            font_size_pt=13, bold=True, color=GOLD)

    # ── Block 4 ──
    add_paragraph(rt, "",  font_size_pt=6, color=WHITE)   # spacer
    add_paragraph(rt, "◆  实盘摩擦 / Live Trading Friction",
                  font_size_pt=11.5, bold=True, color=AMBER_LIGHT, space_before_pt=4)
    add_paragraph(rt,
                  "当前模型按0.1%/每边统一建模，未包含期货展期成本与隔夜资金成本，实盘收益预计低于回测约10-20%。",
                  font_size_pt=10, color=WHITE, space_before_pt=2)
    add_paragraph(rt,
                  "Current model uses 0.1%/side uniform friction. Futures rollover costs and overnight funding not yet modeled — live returns estimated 10-20% below backtest.",
                  font_size_pt=9.5, italic=True, color=LIGHT_GRAY, space_before_pt=1)

    # ── Vertical divider between columns ─────────────────────────────────────
    divider = slide.shapes.add_shape(
        1,
        left=Emu(6050000), top=COL_TOP,
        width=Emu(12700), height=COL_HEIGHT
    )
    add_solid_fill(divider, DARK_GOLD)
    divider.line.fill.background()

    # ── Bottom disclaimer ─────────────────────────────────────────────────────
    disc_top = Emu(5100000)
    disc_box = add_text_box(slide,
                            left=Emu(548640), top=disc_top,
                            width=Emu(10500000), height=Emu(600000))
    dt = disc_box.text_frame
    dt.word_wrap = True

    # Separator line above disclaimer
    disc_sep = slide.shapes.add_shape(
        1,
        left=Emu(548640), top=Emu(5050000),
        width=Emu(11064240), height=Emu(6350)
    )
    add_solid_fill(disc_sep, DARK_GOLD)
    disc_sep.line.fill.background()

    p = dt.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r,
            "All backtests use walk-forward methodology with no look-ahead bias. Past performance does not guarantee future results.",
            font_size_pt=8, color=DIM_GRAY, italic=True)

    add_paragraph(dt,
                  "所有回测采用滚动验证，无未来函数。历史表现不代表未来收益。",
                  font_size_pt=8, color=DIM_GRAY, italic=True, space_before_pt=1)

    # ── Page number (bottom-right) ────────────────────────────────────────────
    pg_box = add_text_box(slide,
                          left=Emu(11094415), top=Emu(6492240),
                          width=Emu(914400), height=Emu(274320))
    pg_tf = pg_box.text_frame
    p = pg_tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, "11 / 16", font_size_pt=8, color=DIM_GRAY)

    return slide


def move_slide_to_index(prs: Presentation, slide, target_index: int):
    """
    Move a slide (already appended) to the given target_index (0-based)
    by manipulating the presentation XML.
    """
    xml_slides = prs.slides._sldIdLst
    # The new slide is currently last — get its element
    slides_list = list(xml_slides)
    new_slide_elem = slides_list[-1]
    # Remove from current position
    xml_slides.remove(new_slide_elem)
    # Insert at target position
    xml_slides.insert(target_index, new_slide_elem)


def main():
    prs = Presentation(PPTX_PATH)
    print(f"Original slide count: {len(prs.slides)}")

    # Build new slide (appended at end)
    build_slide(prs)
    print(f"After append: {len(prs.slides)} slides")

    # Move it to index 10 (11th position, after current slide 10)
    new_slide = prs.slides[-1]
    move_slide_to_index(prs, new_slide, target_index=10)

    # Verify order
    print("Slide order after insertion:")
    for i, s in enumerate(prs.slides):
        title_text = ""
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    title_text = t[:70]
                    break
        print(f"  {i+1:2d}. {title_text}")

    prs.save(PPTX_PATH)
    print(f"\nSaved. Final slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
