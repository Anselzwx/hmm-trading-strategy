"""
HMM Strategy Pitch Deck Generator
Generates a high-end investor-style PPTX presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ─────────────────────────────────────────────────────────────
BG      = RGBColor(15, 17, 23)       # #0F1117
GOLD    = RGBColor(255, 215, 64)     # #FFD740
GREEN   = RGBColor(0, 230, 118)      # #00E676
BLUE    = RGBColor(96, 165, 250)     # #60A5FA
PURPLE  = RGBColor(167, 139, 250)    # #A78BFA
WHITE   = RGBColor(255, 255, 255)
LGRAY   = RGBColor(180, 180, 180)
DGRAY   = RGBColor(100, 100, 100)
DARKBG2 = RGBColor(22, 25, 35)       # slightly lighter for cards

# ── Dimensions (16:9 widescreen) ──────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri",
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_gold_top_bar(slide, height_pt=4):
    """Full-width gold line at the very top of the slide."""
    add_rect(slide,
             left=0, top=0,
             width=SLIDE_W, height=Pt(height_pt),
             fill_color=GOLD)

def add_page_number(slide, page_num, total=15):
    txt = f"{page_num} / {total}"
    add_text_box(slide, txt,
                 left=SLIDE_W - Inches(1.2), top=SLIDE_H - Inches(0.4),
                 width=Inches(1.0), height=Inches(0.3),
                 font_size=10, color=DGRAY, align=PP_ALIGN.RIGHT)

def add_title(slide, title_en, title_zh=None,
              top=Inches(0.5), left=Inches(0.6), width=Inches(12.1)):
    """Add a bilingual title block."""
    h = Inches(0.7) if not title_zh else Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_en
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = GOLD
    r.font.name = "Calibri"

    if title_zh:
        from pptx.util import Pt as P
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = title_zh
        r2.font.size = Pt(16)
        r2.font.bold = False
        r2.font.color.rgb = LGRAY
        r2.font.name = "Microsoft YaHei"

def add_bullet(slide, items, left, top, width, height,
               bullet_color=GOLD, text_color=WHITE, font_size=16,
               line_spacing_pt=8):
    """Add a list of bullet items (list of strings) to a textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        r_bullet = p.add_run()
        r_bullet.text = "▸  "
        r_bullet.font.size = Pt(font_size)
        r_bullet.font.color.rgb = bullet_color
        r_bullet.font.name = "Calibri"
        # text
        r = p.add_run()
        r.text = item
        r.font.size = Pt(font_size)
        r.font.color.rgb = text_color
        r.font.name = "Calibri"

        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(line_spacing_pt * 100))

    return txBox

def add_divider(slide, top, color=GOLD, alpha_fraction=0.3):
    """Thin horizontal divider line."""
    add_rect(slide,
             left=Inches(0.6), top=top,
             width=Inches(12.1), height=Pt(1),
             fill_color=RGBColor(80, 65, 10))

def add_card(slide, left, top, width, height,
             fill=DARKBG2, border_color=GOLD, border_width_pt=1.5):
    shape = add_rect(slide, left, top, width, height,
                     fill_color=fill,
                     line_color=border_color,
                     line_width=Pt(border_width_pt))
    return shape

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide, height_pt=6)

    # Left accent bar
    add_rect(slide, left=0, top=Inches(1.0),
             width=Pt(6), height=Inches(5.5),
             fill_color=GOLD)

    # Big title
    add_text_box(slide,
                 "HMM Regime-Based\nTrading System",
                 left=Inches(0.7), top=Inches(1.4),
                 width=Inches(11.0), height=Inches(2.0),
                 font_size=48, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide,
                 "基于隐马尔可夫模型的多资产量化交易系统",
                 left=Inches(0.7), top=Inches(3.5),
                 width=Inches(11.0), height=Inches(0.6),
                 font_size=22, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

    # Decorative gold line under subtitle
    add_rect(slide, left=Inches(0.7), top=Inches(4.2),
             width=Inches(5.0), height=Pt(2),
             fill_color=GOLD)

    # Bottom info
    add_text_box(slide,
                 "LILYN AI Quant Strategy  ·  2026",
                 left=Inches(0.7), top=Inches(6.5),
                 width=Inches(11.0), height=Inches(0.5),
                 font_size=14, color=LGRAY, align=PP_ALIGN.LEFT)

    add_page_number(slide, 1)


def slide_02_problem(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide,
              "Why Traditional Investing Fails",
              "传统投资的局限")
    add_divider(slide, top=Inches(1.65))

    bullets = [
        "Markets shift between regimes — bull, bear, crash. Static strategies can't adapt.\n    市场在不同状态间切换，静态策略无法适应",
        "Buy & Hold suffers 30-50% drawdowns that take years to recover\n    被动持有面临30-50%回撤，需数年恢复",
        "Human emotion drives poor timing decisions\n    人类情绪导致错误的择时决策",
    ]
    add_bullet(slide, bullets,
               left=Inches(0.8), top=Inches(1.9),
               width=Inches(11.5), height=Inches(4.5),
               font_size=17, line_spacing_pt=18)

    add_page_number(slide, 2)


def slide_03_solution(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide,
              "HMM Regime Detection",
              "隐马尔可夫模型状态识别")
    add_divider(slide, top=Inches(1.65))

    bullets = [
        "Gaussian Hidden Markov Model identifies 7 distinct market regimes\n    高斯HMM识别7种市场状态",
        "Walk-Forward validation ensures no look-ahead bias\n    滚动验证消除未来函数",
        "14-signal confirmation system filters false positives\n    14信号确认系统过滤假信号",
    ]
    add_bullet(slide, bullets,
               left=Inches(0.8), top=Inches(1.9),
               width=Inches(11.5), height=Inches(4.5),
               font_size=17, line_spacing_pt=18)

    add_page_number(slide, 3)


def slide_04_architecture(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "System Architecture", "系统架构")
    add_divider(slide, top=Inches(1.65))

    steps = [
        ("01", "Data Ingestion\n数据摄取", "Daily OHLCV via Yahoo Finance", BLUE),
        ("02", "HMM State Detection\n状态识别", "7-state Gaussian HMM Walk-Forward", GOLD),
        ("03", "Signal Confirmation\n信号确认", "14 Technical Indicators Voting", GREEN),
        ("04", "Execution & Risk\n执行与风控", "Stop-Loss + Position Sizing + Email Alert", PURPLE),
    ]

    card_w = Inches(2.9)
    card_h = Inches(3.8)
    gap = Inches(0.3)
    start_left = Inches(0.5)
    top = Inches(2.0)

    for i, (num, step_title, desc, color) in enumerate(steps):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=color, border_width_pt=2)

        # Number badge
        add_text_box(slide, num,
                     left=lft + Inches(0.1), top=top + Inches(0.15),
                     width=Inches(0.7), height=Inches(0.5),
                     font_size=28, bold=True, color=color)

        # Step title
        add_text_box(slide, step_title,
                     left=lft + Inches(0.15), top=top + Inches(0.75),
                     width=card_w - Inches(0.3), height=Inches(1.2),
                     font_size=15, bold=True, color=WHITE,
                     font_name="Microsoft YaHei")

        # Arrow connector (except last)
        if i < 3:
            arrow_left = lft + card_w + gap * 0.1
            add_text_box(slide, "→",
                         left=arrow_left, top=top + card_h / 2 - Inches(0.2),
                         width=gap * 0.8, height=Inches(0.4),
                         font_size=20, color=GOLD, align=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, desc,
                     left=lft + Inches(0.15), top=top + Inches(2.1),
                     width=card_w - Inches(0.3), height=Inches(1.4),
                     font_size=13, color=LGRAY)

    add_page_number(slide, 4)


def slide_05_assets(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Multi-Asset Coverage", "多资产覆盖")
    add_divider(slide, top=Inches(1.65))

    assets = [
        ("Gold Futures  GC=F", "黄金期货\n全球避险核心资产", GOLD),
        ("Silver Futures  SI=F", "白银期货\n工业+贵金属双重属性", LGRAY),
        ("Apple Inc.  AAPL", "科技龙头\n流动性最佳", BLUE),
    ]

    card_w = Inches(3.8)
    card_h = Inches(3.6)
    gap = Inches(0.45)
    start_left = Inches(0.55)
    top = Inches(2.1)

    for i, (name, zh, color) in enumerate(assets):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=color, border_width_pt=2)

        # Top color bar
        add_rect(slide, left=lft, top=top,
                 width=card_w, height=Pt(5),
                 fill_color=color)

        add_text_box(slide, name,
                     left=lft + Inches(0.2), top=top + Inches(0.3),
                     width=card_w - Inches(0.4), height=Inches(0.8),
                     font_size=20, bold=True, color=color)

        add_text_box(slide, zh,
                     left=lft + Inches(0.2), top=top + Inches(1.3),
                     width=card_w - Inches(0.4), height=Inches(1.5),
                     font_size=16, color=WHITE,
                     font_name="Microsoft YaHei")

    add_page_number(slide, 5)


def _perf_slide(prs, page_num, title_en, title_zh, metrics, note=None, accent=GOLD):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, title_en, title_zh)
    add_divider(slide, top=Inches(1.65))

    # Big highlight metrics (top row)
    highlight = metrics[:3]
    rest = metrics[3:]

    card_w = Inches(3.8)
    card_h = Inches(1.6)
    gap = Inches(0.4)
    top = Inches(1.85)
    start_left = Inches(0.55)

    for i, (label, value) in enumerate(highlight):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=accent, border_width_pt=1.5)
        add_text_box(slide, value,
                     left=lft + Inches(0.15), top=top + Inches(0.1),
                     width=card_w - Inches(0.3), height=Inches(0.9),
                     font_size=30, bold=True, color=accent,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     left=lft + Inches(0.15), top=top + Inches(1.0),
                     width=card_w - Inches(0.3), height=Inches(0.45),
                     font_size=13, color=LGRAY, align=PP_ALIGN.CENTER)

    # Secondary metrics
    sec_top = top + card_h + Inches(0.3)
    sec_card_w = Inches(2.0)
    sec_card_h = Inches(1.3)
    sec_gap = Inches(0.22)
    n = len(rest)
    total_w = n * sec_card_w + (n - 1) * sec_gap
    sec_start = (SLIDE_W - total_w) / 2

    for i, (label, value) in enumerate(rest):
        lft = sec_start + i * (sec_card_w + sec_gap)
        add_card(slide, lft, sec_top, sec_card_w, sec_card_h,
                 border_color=GREEN, border_width_pt=1)
        add_text_box(slide, value,
                     left=lft + Inches(0.1), top=sec_top + Inches(0.05),
                     width=sec_card_w - Inches(0.2), height=Inches(0.75),
                     font_size=22, bold=True, color=GREEN,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     left=lft + Inches(0.1), top=sec_top + Inches(0.8),
                     width=sec_card_w - Inches(0.2), height=Inches(0.4),
                     font_size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    if note:
        add_text_box(slide, note,
                     left=Inches(0.6), top=Inches(6.7),
                     width=Inches(12.0), height=Inches(0.4),
                     font_size=11, color=DGRAY)

    add_page_number(slide, page_num)


def slide_06_gold(prs):
    metrics = [
        ("Total Return", "+4,594%"),
        ("Annual Return", "43.5%"),
        ("Sharpe Ratio", "1.35"),
        ("Max Drawdown", "-9.9%"),
        ("Sortino", "9.91"),
        ("Calmar", "4.38"),
        ("Win Rate", "63.6%"),
        ("Profit Factor", "4.33"),
        ("Trades (10Y)", "33"),
    ]
    _perf_slide(prs, 6,
                "Gold Strategy Performance",
                "黄金策略表现",
                metrics,
                note="vs Buy & Hold: +331%  |  10-year backtest 2016–2026",
                accent=GOLD)


def slide_07_silver(prs):
    metrics = [
        ("Total Return", "+925%"),
        ("Annual Return", "24.4%"),
        ("Sharpe Ratio", "0.96"),
        ("Max Drawdown", "-14.5%"),
        ("Sortino", "4.45"),
        ("Calmar", "1.68"),
        ("Win Rate", "50.0%"),
        ("Profit Factor", "3.16"),
        ("vs B&H", "+444%"),
    ]
    _perf_slide(prs, 7,
                "Silver Strategy Performance",
                "白银策略表现",
                metrics,
                note="vs Buy & Hold: +444%  |  10-year backtest 2016–2026",
                accent=LGRAY)


def slide_08_aapl(prs):
    metrics = [
        ("Total Return", "+189%"),
        ("Annual Return", "10.9%"),
        ("Sharpe Ratio", "0.97"),
        ("Max Drawdown", "-16.9%"),
        ("Win Rate", "62.9%"),
        ("Profit Factor", "3.75"),
    ]
    _perf_slide(prs, 8,
                "Apple Strategy Performance",
                "苹果策略表现",
                metrics,
                note="AAPL B&H: +1004%  |  Strategy focuses on risk control, not trend-chasing  /  策略侧重风控而非追涨",
                accent=BLUE)


def slide_09_risk(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Risk Control Framework", "风险控制框架")
    add_divider(slide, top=Inches(1.65))

    modules = [
        ("Hard Stop-Loss\n硬止损",
         "Gold/Silver −8%  |  AAPL −6%\nTriggered immediately on breach\n触价即出", GOLD),
        ("HMM Bear Filter\nHMM熊市过滤",
         "Bear probability >60% → alert\nBear state → immediate exit\nbear state直接出场", GREEN),
        ("Position Sizing\n仓位管理",
         "Signal score → position weight\n40% – 100% capital allocation\n信号分数驱动仓位比例", BLUE),
        ("Regime Reduce\n状态缩减",
         "Consecutive bear signals:\n100% → 50% → 25%\n连续bear信号逐步减仓", PURPLE),
    ]

    card_w = Inches(2.9)
    card_h = Inches(3.6)
    gap = Inches(0.3)
    start_left = Inches(0.5)
    top = Inches(1.9)

    for i, (title, body, color) in enumerate(modules):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=color, border_width_pt=2)
        add_rect(slide, left=lft, top=top,
                 width=card_w, height=Pt(4),
                 fill_color=color)
        add_text_box(slide, title,
                     left=lft + Inches(0.15), top=top + Inches(0.15),
                     width=card_w - Inches(0.3), height=Inches(1.1),
                     font_size=15, bold=True, color=color,
                     font_name="Microsoft YaHei")
        add_text_box(slide, body,
                     left=lft + Inches(0.15), top=top + Inches(1.3),
                     width=card_w - Inches(0.3), height=Inches(2.0),
                     font_size=13, color=LGRAY)

    add_page_number(slide, 9)


def slide_10_wf(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Walk-Forward Backtesting", "滚动验证回测")
    add_divider(slide, top=Inches(1.65))

    bullets = [
        "60% training / 40% rolling test window  |  60%训练窗口，40%滚动测试",
        "No look-ahead bias — each signal uses only past data  |  零未来函数",
        "Subsample validation: 2016–2020 / 2021–2023 / 2024–2026 all positive  |  三段子样本全部正收益",
        "10-year backtest: 2016–2026  |  10年全样本回测",
    ]
    add_bullet(slide, bullets,
               left=Inches(0.8), top=Inches(1.9),
               width=Inches(11.5), height=Inches(3.5),
               font_size=17, line_spacing_pt=20)

    # Mini timeline bar
    bar_top = Inches(5.5)
    bar_h = Inches(0.45)
    bar_total_w = Inches(11.2)
    bar_left = Inches(0.9)

    segments = [
        ("Training 2016–2019", 0.4, GOLD),
        ("Val 2020–2021", 0.2, BLUE),
        ("Training 2020–2022", 0.4, GOLD),
        ("Val 2023–2024", 0.2, GREEN),
        ("OOS 2024–2026", 0.25, PURPLE),
    ]
    total_parts = sum(s[1] for s in segments)
    cursor = bar_left
    for label, weight, color in segments:
        w = bar_total_w * (weight / total_parts)
        add_rect(slide, cursor, bar_top, w, bar_h, fill_color=color)
        add_text_box(slide, label,
                     left=cursor, top=bar_top + bar_h + Pt(4),
                     width=w, height=Inches(0.35),
                     font_size=9, color=color, align=PP_ALIGN.CENTER)
        cursor += w

    add_page_number(slide, 10)


def slide_11_automation(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Fully Automated Pipeline", "全自动化流程")
    add_divider(slide, top=Inches(1.65))

    steps = [
        ("17:00 EST", "Daily Close triggers\nsignal_generator.py", GOLD),
        ("Inference", "HMM model inference\n+ 14-signal scoring", BLUE),
        ("Alert", "monitor.py sends\nemail on ENTER/EXIT", GREEN),
        ("Dashboard", "hmm-trading.streamlit.app\nupdates live", PURPLE),
    ]

    card_w = Inches(2.7)
    card_h = Inches(2.8)
    gap = Inches(0.45)
    start_left = Inches(0.7)
    top = Inches(2.2)

    for i, (time_label, desc, color) in enumerate(steps):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=color, border_width_pt=2)
        add_text_box(slide, time_label,
                     left=lft + Inches(0.15), top=top + Inches(0.15),
                     width=card_w - Inches(0.3), height=Inches(0.6),
                     font_size=20, bold=True, color=color,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     left=lft + Inches(0.15), top=top + Inches(0.85),
                     width=card_w - Inches(0.3), height=Inches(1.6),
                     font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

        if i < 3:
            arr_lft = lft + card_w + gap * 0.05
            add_text_box(slide, "→",
                         left=arr_lft, top=top + card_h / 2 - Inches(0.2),
                         width=gap * 0.9, height=Inches(0.4),
                         font_size=22, color=GOLD, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "Zero manual intervention required  /  零人工干预",
                 left=Inches(0.7), top=Inches(5.6),
                 width=Inches(12.0), height=Inches(0.5),
                 font_size=14, color=LGRAY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 11)


def slide_12_dashboard(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Real-Time Dashboard", "实时监控面板")
    add_divider(slide, top=Inches(1.65))

    tabs = [
        ("Today's Signals\n今日信号", "Daily action recommendations", GOLD),
        ("Portfolio\n组合", "Equal-weight portfolio performance", BLUE),
        ("Gold\n黄金", "Per-asset deep dive analytics", GOLD),
        ("Silver\n白银", "Per-asset deep dive analytics", LGRAY),
        ("AAPL\n苹果", "Per-asset deep dive analytics", BLUE),
    ]

    card_w = Inches(2.3)
    card_h = Inches(3.2)
    gap = Inches(0.25)
    start_left = Inches(0.55)
    top = Inches(2.0)

    for i, (name, desc, color) in enumerate(tabs):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=color, border_width_pt=1.5)
        add_rect(slide, left=lft, top=top,
                 width=card_w, height=Pt(4), fill_color=color)
        add_text_box(slide, name,
                     left=lft + Inches(0.1), top=top + Inches(0.15),
                     width=card_w - Inches(0.2), height=Inches(1.2),
                     font_size=14, bold=True, color=color,
                     align=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
        add_text_box(slide, desc,
                     left=lft + Inches(0.1), top=top + Inches(1.5),
                     width=card_w - Inches(0.2), height=Inches(1.4),
                     font_size=12, color=LGRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "http://hmm-trading.streamlit.app",
                 left=Inches(0.6), top=Inches(5.8),
                 width=Inches(12.0), height=Inches(0.4),
                 font_size=14, color=BLUE, align=PP_ALIGN.CENTER)

    add_page_number(slide, 12)


def slide_13_edge(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Why This System Wins", "核心竞争优势")
    add_divider(slide, top=Inches(1.65))

    # Table
    headers = ["", "HMM Strategy", "Traditional Quant", "Buy & Hold"]
    rows = [
        ["Regime Awareness\n状态感知", "✅ Yes", "❌ No", "❌ No"],
        ["Max Drawdown\n最大回撤", "-9.9%", "-20~40%", "-30~50%"],
        ["Automation\n自动化", "✅ Full", "Partial", "✅ Passive"],
        ["Sharpe Ratio", "1.35", "0.5~1.0", "0.4~0.7"],
    ]

    col_widths = [Inches(2.6), Inches(2.9), Inches(3.1), Inches(2.9)]
    row_h = Inches(0.78)
    tbl_left = Inches(0.6)
    tbl_top = Inches(1.85)
    header_colors = [BG, GOLD, BLUE, LGRAY]

    # Header row
    cursor = tbl_left
    for j, (hdr, w, hc) in enumerate(zip(headers, col_widths, header_colors)):
        add_rect(slide, cursor, tbl_top, w, row_h,
                 fill_color=hc if j > 0 else DARKBG2,
                 line_color=BG, line_width=Pt(1))
        txt_color = BG if j > 0 else GOLD
        add_text_box(slide, hdr,
                     left=cursor + Inches(0.1), top=tbl_top + Inches(0.1),
                     width=w - Inches(0.2), height=row_h - Inches(0.1),
                     font_size=14, bold=True, color=txt_color,
                     align=PP_ALIGN.CENTER)
        cursor += w

    # Data rows
    row_colors = [DARKBG2, RGBColor(20, 22, 32)]
    for ri, row in enumerate(rows):
        row_top = tbl_top + (ri + 1) * row_h
        cursor = tbl_left
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            bg = row_colors[ri % 2]
            add_rect(slide, cursor, row_top, w, row_h,
                     fill_color=bg, line_color=DGRAY, line_width=Pt(0.5))
            # Highlight HMM column in gold, first col in white
            cell_color = GOLD if j == 1 else (LGRAY if j == 0 else WHITE)
            if j == 1 and "✅" in cell:
                cell_color = GREEN
            add_text_box(slide, cell,
                         left=cursor + Inches(0.1), top=row_top + Inches(0.1),
                         width=w - Inches(0.2), height=row_h - Inches(0.15),
                         font_size=14, bold=(j == 1), color=cell_color,
                         align=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
            cursor += w

    add_page_number(slide, 13)


def slide_14_roadmap(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide)
    add_title(slide, "Next Steps", "下一步计划")
    add_divider(slide, top=Inches(1.65))

    phases = [
        ("Q2 2026", "Live trading via IBKR API\nIBKR API实盘对接", GOLD),
        ("Q3 2026", "Asset universe expansion\n(ETFs, Crypto)\n资产池扩展", BLUE),
        ("Q4 2026", "NLP sentiment layer\nfrom news data\n新闻情绪因子接入", GREEN),
        ("2027", "Multi-strategy ensemble\n多策略组合", PURPLE),
    ]

    card_w = Inches(2.85)
    card_h = Inches(3.5)
    gap = Inches(0.35)
    start_left = Inches(0.6)
    top = Inches(2.0)

    for i, (period, desc, color) in enumerate(phases):
        lft = start_left + i * (card_w + gap)
        add_card(slide, lft, top, card_w, card_h,
                 border_color=color, border_width_pt=2)
        add_rect(slide, lft, top, card_w, Pt(5), fill_color=color)

        add_text_box(slide, period,
                     left=lft + Inches(0.15), top=top + Inches(0.15),
                     width=card_w - Inches(0.3), height=Inches(0.65),
                     font_size=22, bold=True, color=color,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, desc,
                     left=lft + Inches(0.15), top=top + Inches(1.0),
                     width=card_w - Inches(0.3), height=Inches(2.2),
                     font_size=14, color=WHITE, align=PP_ALIGN.CENTER,
                     font_name="Microsoft YaHei")

    add_page_number(slide, 14)


def slide_15_contact(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG)
    add_gold_top_bar(slide, height_pt=6)

    # Left accent
    add_rect(slide, left=0, top=Inches(1.0),
             width=Pt(6), height=Inches(5.5),
             fill_color=GOLD)

    # Thank You
    add_text_box(slide, "Thank You",
                 left=Inches(0.7), top=Inches(1.2),
                 width=Inches(11.0), height=Inches(0.9),
                 font_size=44, bold=True, color=GOLD,
                 align=PP_ALIGN.LEFT)

    add_text_box(slide, "谢谢",
                 left=Inches(0.7), top=Inches(2.1),
                 width=Inches(11.0), height=Inches(0.7),
                 font_size=28, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Microsoft YaHei")

    add_rect(slide, left=Inches(0.7), top=Inches(2.9),
             width=Inches(4.0), height=Pt(2),
             fill_color=GOLD)

    info_lines = [
        ("LILYN AI Quant Strategy", WHITE, 18, True),
        ("Contact:  anselwilliam789@gmail.com", LGRAY, 15, False),
        ("Dashboard:  hmm-trading.streamlit.app", BLUE, 15, False),
        ("GitHub:  github.com/Anselzwx/hmm-trading-strategy", LGRAY, 14, False),
    ]

    y = Inches(3.1)
    for text, color, fsize, bold in info_lines:
        add_text_box(slide, text,
                     left=Inches(0.7), top=y,
                     width=Inches(11.5), height=Inches(0.5),
                     font_size=fsize, bold=bold, color=color,
                     align=PP_ALIGN.LEFT)
        y += Inches(0.52)

    # Disclaimer
    add_text_box(slide,
                 "Past performance does not guarantee future results. For research purposes only.\n"
                 "历史表现不代表未来收益，仅供研究参考。",
                 left=Inches(0.7), top=Inches(6.4),
                 width=Inches(12.0), height=Inches(0.7),
                 font_size=10, color=DGRAY,
                 font_name="Microsoft YaHei")

    add_page_number(slide, 15)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_architecture(prs)
    slide_05_assets(prs)
    slide_06_gold(prs)
    slide_07_silver(prs)
    slide_08_aapl(prs)
    slide_09_risk(prs)
    slide_10_wf(prs)
    slide_11_automation(prs)
    slide_12_dashboard(prs)
    slide_13_edge(prs)
    slide_14_roadmap(prs)
    slide_15_contact(prs)

    out_path = "/Users/zhaowenxuan/Desktop/工作/HMM/HMM_Strategy_Pitch.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
