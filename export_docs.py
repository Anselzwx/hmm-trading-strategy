"""
export_docs.py
--------------
生成：
  1. HMM_Strategy_Report.docx  —— 完整策略说明文档
  2. HMM_Strategy_Deck.pptx    —— 简约高级 PPT（8 页）
"""

import os
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PPTpt, Emu
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PPTpt

OUT_DIR = os.path.dirname(__file__)

# ──────────────────────────────────────────────────────────────
# 颜色常量
# ──────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x0A, 0x0A, 0x0A)
DARK_GRAY  = RGBColor(0x1E, 0x1E, 0x2E)
MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GREEN      = RGBColor(0x00, 0xC8, 0x6E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x60, 0xA5, 0xFA)

PPT_BG     = PPTColor(0x0A, 0x0E, 0x1A)
PPT_WHITE  = PPTColor(0xF1, 0xF5, 0xF9)
PPT_GREEN  = PPTColor(0x00, 0xC8, 0x6E)
PPT_GRAY   = PPTColor(0x64, 0x74, 0x8B)
PPT_ACCENT = PPTColor(0x60, 0xA5, 0xFA)
PPT_CARD   = PPTColor(0x14, 0x1C, 0x2E)


# ══════════════════════════════════════════════════════════════
# WORD 文档
# ══════════════════════════════════════════════════════════════

def set_cell_bg(cell, hex_color: str):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  hex_color)
    tcPr.append(shd)


def add_heading(doc, text, level=1, color=None):
    p    = doc.add_paragraph()
    run  = p.add_run(text)
    size = {1: 22, 2: 16, 3: 13}.get(level, 13)
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = color or BLACK
    p.paragraph_format.space_before = Pt(16 if level == 1 else 10)
    p.paragraph_format.space_after  = Pt(6)
    return p


def add_body(doc, text, color=None, size=11):
    p   = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size      = Pt(size)
    run.font.color.rgb = color or DARK_GRAY
    p.paragraph_format.space_after = Pt(4)
    return p


def add_bullet(doc, text, level=0):
    p   = doc.add_paragraph(style='List Bullet')
    run = p.add_run(text)
    run.font.size      = Pt(10.5)
    run.font.color.rgb = DARK_GRAY
    p.paragraph_format.left_indent   = Cm(0.5 + level * 0.8)
    p.paragraph_format.space_after   = Pt(3)
    return p


def add_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Table Grid'

    # 表头
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        set_cell_bg(cell, '1E2130')
        run  = cell.paragraphs[0].add_run(h)
        run.font.bold       = True
        run.font.color.rgb  = WHITE
        run.font.size       = Pt(10)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 数据行
    for ri, row in enumerate(rows):
        tr = table.rows[ri + 1]
        bg = 'F8FAFC' if ri % 2 == 0 else 'EFF6FF'
        for ci, val in enumerate(row):
            cell = tr.cells[ci]
            set_cell_bg(cell, bg)
            run  = cell.paragraphs[0].add_run(str(val))
            run.font.size      = Pt(10)
            run.font.color.rgb = DARK_GRAY
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    if col_widths:
        for i, w in enumerate(col_widths):
            for row in table.rows:
                row.cells[i].width = Cm(w)

    doc.add_paragraph()
    return table


def build_word():
    doc = Document()

    # 页面边距
    for sec in doc.sections:
        sec.top_margin    = Cm(2.5)
        sec.bottom_margin = Cm(2.5)
        sec.left_margin   = Cm(3.0)
        sec.right_margin  = Cm(3.0)

    # 封面
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('HMM Regime-Based Trading Strategy')
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(0x0A, 0x0E, 0x1A)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run('Phase 1 策略说明文档')
    r2.font.size      = Pt(14)
    r2.font.color.rgb = MID_GRAY

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = p3.add_run('资产：AAPL · Gold (GC=F) · Silver (SI=F)　　数据截至：2026-04-16')
    r3.font.size      = Pt(11)
    r3.font.color.rgb = MID_GRAY

    doc.add_paragraph()
    doc.add_paragraph('─' * 60)
    doc.add_paragraph()

    # 1. 项目概述
    add_heading(doc, '一、项目概述', 1)
    add_body(doc, '本项目构建了一套基于隐马尔可夫模型（Gaussian HMM）的量化交易策略回测系统，'
                  '研究对象为三个资产：苹果股票（AAPL）、黄金期货（Gold）、白银期货（Silver）。')
    add_body(doc, '核心研究问题：')
    add_bullet(doc, '用机器学习自动识别市场状态（牛市 / 熊市 / 震荡）')
    add_bullet(doc, '只在"天时地利人和"均满足时入场，避免随机交易')
    add_bullet(doc, '在真实风控约束下，策略能否跑赢买入持有基准？')

    # 2. HMM 引擎
    add_heading(doc, '二、市场状态识别（HMM 引擎）', 1)
    add_body(doc, '模型：7 状态 Gaussian HMM，训练特征如下：')
    add_table(doc,
        ['特征', '计算方式', '含义'],
        [
            ['对数收益率', 'ln(Close_t / Close_{t-1}) × 100', '当前 K 线的涨跌幅'],
            ['价格区间比', '(High - Low) / Close × 100', '衡量 K 线内部波动强度'],
            ['成交量波动率', 'rolling std(log Volume, 24)', '反映市场情绪变化'],
        ],
        col_widths=[3.5, 6, 5]
    )
    add_body(doc, '状态标记规则：')
    add_bullet(doc, 'Bull Run：收益率均值最高的状态（最强牛市）')
    add_bullet(doc, 'Bull+：收益率第二高的状态（次级牛市）')
    add_bullet(doc, 'Bear/Crash：收益率最低的状态（熊市 / 崩盘）')
    add_bullet(doc, 'Neutral-X：其余 4 个状态（震荡 / 过渡）')

    # 3. Walk-Forward
    add_heading(doc, '三、Walk-Forward 滚动训练（消除数据泄露）', 1)
    add_body(doc, '传统回测的致命缺陷：用全量数据训练模型后再用同一段数据验证，'
                  '相当于"用答案训练模型再考同一套题"，统计学上称为 Look-ahead Bias（未来数据泄露）。')
    add_body(doc, 'Phase 1 采用 Walk-Forward 滚动训练方案：')
    add_bullet(doc, '用前 60% 数据训练 HMM，预测之后的数据')
    add_bullet(doc, '每隔 10% 向前滚动，重新训练一次')
    add_bullet(doc, '每个预测点使用的模型，只见过该时刻之前的数据')
    add_bullet(doc, '这才是真实部署环境下策略能跑出的结果')

    # 4. 14 信号
    add_heading(doc, '四、入场过滤（14 信号投票系统）', 1)
    add_body(doc, 'HMM 判断为牛市状态后，还需满足 14 个技术指标中的至少 9 个，才真正入场。')
    add_table(doc,
        ['#', '信号', '条件', '含义'],
        [
            ['1',  'RSI',          'RSI < 90',              '没有严重超买'],
            ['2',  '动量',         'Momentum > 1%',         '价格在加速上涨'],
            ['3',  '波动率',       'Volatility < 6%',       '市场不过于剧烈'],
            ['4',  '成交量',       'Volume > SMA20',        '有真实买盘支撑'],
            ['5',  'ADX',          'ADX > 25',              '趋势足够强'],
            ['6',  'EMA50',        'Price > EMA50',         '中期趋势向上'],
            ['7',  'EMA200',       'Price > EMA200',        '长期趋势向上'],
            ['8',  'MACD',         'MACD > Signal',         '短期动能看多'],
            ['9',  'Bollinger',    'Price > BB 中轨',       '处于布林带多头区'],
            ['10', 'Stochastic',   '%K 上穿 %D 且 < 80',   '动量回升未超买'],
            ['11', 'Williams %R',  '%R < -20',              '未进极度超买区'],
            ['12', 'CCI',          'CCI > 0',               '价格高于统计均值'],
            ['13', 'OBV',          'OBV > OBV EMA',         '量能支撑上涨'],
            ['14', '高点距离',     '距高点回撤 < 30%',      '不在深跌中接刀'],
        ],
        col_widths=[1, 3, 4.5, 4.5]
    )

    # 5. 风控
    add_heading(doc, '五、风险管理规则', 1)
    add_table(doc,
        ['规则', '参数', '逻辑'],
        [
            ['固定止损',   '-8% / 笔',        '单笔亏损超过 8% 立即出场，不等 HMM 反应'],
            ['状态止损',   'Bear/Crash',       'HMM 翻熊立即出场'],
            ['冷静期',     '48h / 2 交易日',   '出场后强制等待，避免反复被割'],
            ['最长持仓',   '30天 / 60交易日',  '到期自动止盈再重新评估'],
            ['杠杆',       '2.5×',            '放大盈亏，同步放大风险'],
        ],
        col_widths=[3, 4, 7]
    )

    add_heading(doc, '仓位管理（信号强度决定仓位大小）', 2)
    add_table(doc,
        ['信号得分', '仓位比例', '说明'],
        [
            ['9 / 14',  '40%',  '刚过门槛，保守入场'],
            ['11 / 14', '60%',  '中等信号，适中仓位'],
            ['14 / 14', '100%', '全部满足，全仓入场'],
        ],
        col_widths=[4, 4, 6]
    )

    # 6. 绩效指标
    add_heading(doc, '六、绩效指标说明', 1)
    add_table(doc,
        ['指标', '含义', '参考标准'],
        [
            ['总收益',     '策略整体盈亏百分比',              '越高越好'],
            ['Alpha',      '相对买入持有的超额收益',          '正数才有价值'],
            ['夏普比率',   '每单位风险获得的超额收益',        '> 1 为优秀'],
            ['卡玛比率',   '年化收益 / 最大回撤',            '> 1 为优秀'],
            ['月度胜率',   '盈利月份占总月份的比例',          '> 60% 较稳定'],
            ['最大回撤',   '从峰值到谷值的最大跌幅',          '越小越好'],
            ['盈亏比 R:R', '平均盈利 / 平均亏损',            '> 1.5 较健康'],
            ['vs SPY',     '相对标普 500 的超额收益',        '正数说明跑赢大盘'],
        ],
        col_widths=[3.5, 6, 4.5]
    )

    # 7. 回测结果
    add_heading(doc, '七、当前回测结果（数据截至 2026-04-16）', 1)
    add_table(doc,
        ['资产', '总收益', '夏普比率', '最大回撤', '月度胜率', '交易笔数'],
        [
            ['AAPL',   '+11.9%', '1.04', '-12.1%', '—',   '38'],
            ['Gold',   '+80.6%', '1.51', '-21.5%', '75%', '25'],
            ['Silver', '+80.7%', '0.64', '-46.9%', '50%', '27'],
        ],
        col_widths=[3, 3, 3, 3, 3, 3]
    )
    add_body(doc, '结论分析：')
    add_bullet(doc, 'Gold 表现最优：夏普 1.51，月度胜率 75%，回撤控制在 -21.5%，适合作为核心配置')
    add_bullet(doc, 'AAPL 最为稳健：夏普刚过 1，最大回撤仅 -12.1%，波动最小')
    add_bullet(doc, 'Silver 收益高但波动剧烈：回撤接近 -47%，风险偏高，需谨慎配置')

    # 8. 技术架构
    add_heading(doc, '八、技术架构', 1)
    add_table(doc,
        ['文件', '职责'],
        [
            ['data_loader.py',  '数据拉取（yfinance）、分块处理、本地缓存'],
            ['backtester.py',   'HMM 训练、指标计算、Walk-Forward 回测、绩效统计'],
            ['precompute.py',   '预计算脚本，本地运行后推送结果到云端'],
            ['app.py',          'Streamlit 可视化 Dashboard，秒开'],
            ['results/',        '预计算 pkl 文件，Streamlit Cloud 直接读取'],
        ],
        col_widths=[5, 9]
    )

    # 9. Phase 2
    add_heading(doc, '九、下一步计划（Phase 2）', 1)
    add_bullet(doc, '指标重要性分析：量化 14 个信号的预测贡献度，保留最有效的，去掉冗余的')
    add_bullet(doc, '跨资产相关性：Gold 与 Silver 高度相关，研究是否应该错峰配置')
    add_bullet(doc, '与 SPY 对比：验证策略是否只在牛市有效，还是真正的 Alpha')
    add_bullet(doc, '参数优化：对止损线、持仓上限、信号门槛做网格搜索')

    doc.add_paragraph()
    p_end = doc.add_paragraph()
    p_end.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_end = p_end.add_run('— 仅供学习研究，不构成投资建议 —')
    r_end.font.size      = Pt(10)
    r_end.font.color.rgb = MID_GRAY
    r_end.font.italic    = True

    out = os.path.join(OUT_DIR, 'HMM_Strategy_Report.docx')
    doc.save(out)
    print(f'✅ Word 已保存：{out}')


# ══════════════════════════════════════════════════════════════
# PPT
# ══════════════════════════════════════════════════════════════

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_slide_bg(slide, color: PPTColor):
    bg    = slide.background
    fill  = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = PPTpt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=24, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False):
    txb  = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf   = txb.text_frame
    tf.word_wrap = True
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text           = text
    run.font.size      = PPTpt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color or PPT_WHITE
    return txb


def add_divider(slide, t, color=PPT_GREEN):
    add_rect(slide, 0.6, t, 1.2, 0.04, fill_color=color)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)

    # 左侧绿色竖条
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    # 装饰圆（右下角）
    circle = slide.shapes.add_shape(9, Inches(9.5), Inches(4.5), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PPTColor(0x14, 0x1C, 0x2E)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(9, Inches(10.5), Inches(5.2), Inches(3.5), Inches(3.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = PPTColor(0x00, 0xC8, 0x6E)
    circle2.fill.fore_color.rgb = PPTColor(0x00, 0x40, 0x28)
    circle2.line.fill.background()

    add_text(slide, 'HMM REGIME-BASED', 0.8, 1.6, 9, 1,
             font_size=42, bold=True, color=PPT_WHITE)
    add_text(slide, 'TRADING STRATEGY', 0.8, 2.5, 9, 1,
             font_size=42, bold=True, color=PPT_GREEN)
    add_divider(slide, 3.65)
    add_text(slide, 'Phase 1 策略汇报  ·  AAPL · Gold · Silver', 0.8, 3.8, 10, 0.5,
             font_size=16, color=PPT_GRAY)
    add_text(slide, '数据截至 2026-04-16  ·  Walk-Forward · 14-Signal Voting · 2.5× Leverage',
             0.8, 4.3, 11, 0.5, font_size=12, color=PPT_GRAY)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '01  项目概述', 0.8, 0.4, 10, 0.6,
             font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, '我们在做什么？', 0.8, 0.85, 10, 0.8,
             font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    cards = [
        ('🧠', 'HMM 状态识别', '用机器学习自动识别\n牛市 / 熊市 / 震荡'),
        ('📊', '14 信号投票', '多重技术指标确认\n降低假信号率'),
        ('🛡', '量化风控', '固定止损 + 仓位管理\n+ 冷静期保护'),
        ('✅', 'Walk-Forward', '滚动训练消除\n数据泄露'),
    ]
    for i, (icon, title, desc) in enumerate(cards):
        x = 0.8 + i * 3.1
        add_rect(slide, x, 2.1, 2.8, 3.8, fill_color=PPT_CARD)
        add_text(slide, icon,  x+0.2, 2.3,  2.4, 0.7, font_size=28)
        add_text(slide, title, x+0.2, 3.05, 2.4, 0.6, font_size=14, bold=True, color=PPT_WHITE)
        add_text(slide, desc,  x+0.2, 3.7,  2.4, 1.2, font_size=11, color=PPT_GRAY)

    add_text(slide, '核心问题：只在"天时地利人和"均满足时入场，能否跑赢买入持有？',
             0.8, 6.2, 11.5, 0.6, font_size=13, color=PPT_ACCENT, bold=True)


def slide_hmm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '02  HMM 引擎', 0.8, 0.4, 10, 0.6, font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, '市场状态自动识别', 0.8, 0.85, 10, 0.8, font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    # 左：特征
    add_text(slide, '三个训练特征', 0.8, 2.0, 5.5, 0.5, font_size=14, bold=True, color=PPT_ACCENT)
    feats = [
        ('对数收益率', 'ln(Close_t / Close_{t-1}) × 100', '衡量涨跌幅'),
        ('价格区间比', '(High - Low) / Close × 100',      '衡量波动强度'),
        ('成交量波动率','rolling std(log Vol, 24)',        '反映市场情绪'),
    ]
    for i, (name, formula, desc) in enumerate(feats):
        y = 2.6 + i * 1.1
        add_rect(slide, 0.8, y, 5.6, 0.9, fill_color=PPT_CARD)
        add_text(slide, name,    1.0, y+0.05, 2,   0.4, font_size=12, bold=True, color=PPT_WHITE)
        add_text(slide, formula, 1.0, y+0.45, 3.5, 0.35, font_size=9,  color=PPT_GRAY, italic=True)
        add_text(slide, desc,    4.0, y+0.2,  2.2, 0.5, font_size=11, color=PPT_ACCENT)

    # 右：状态标记
    add_text(slide, '7 个状态的标记规则', 7.0, 2.0, 5.8, 0.5, font_size=14, bold=True, color=PPT_ACCENT)
    states = [
        ('Bull Run',   PPTColor(0x00,0xC8,0x6E), '收益率最高 → 强牛市入场'),
        ('Bull+',      PPTColor(0x00,0x80,0x50), '次高收益率 → 次级牛市'),
        ('Bear/Crash', PPTColor(0xFF,0x52,0x52), '收益率最低 → 立即出场'),
        ('Neutral ×4', PPTColor(0xFF,0xD7,0x40), '其余状态 → 观望'),
    ]
    for i, (name, color, desc) in enumerate(states):
        y = 2.6 + i * 1.05
        add_rect(slide, 7.0, y, 0.18, 0.7, fill_color=color)
        add_rect(slide, 7.22, y, 5.5, 0.7, fill_color=PPT_CARD)
        add_text(slide, name, 7.35, y+0.05, 2, 0.35, font_size=12, bold=True, color=PPT_WHITE)
        add_text(slide, desc, 7.35, y+0.38, 5, 0.3,  font_size=10, color=PPT_GRAY)


def slide_walkforward(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '03  Walk-Forward', 0.8, 0.4, 10, 0.6, font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, '消除数据泄露，让回测可信', 0.8, 0.85, 11, 0.8, font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    # 旧版 vs 新版
    add_rect(slide, 0.8, 2.1, 5.6, 2.4, fill_color=PPTColor(0x2D,0x10,0x10))
    add_text(slide, '❌  旧版（有缺陷）', 1.0, 2.2, 5, 0.5, font_size=13, bold=True, color=PPTColor(0xFF,0x52,0x52))
    add_text(slide, '用全量数据训练 HMM\n再用同一段数据回测\n= 考前把答案给了模型\n= Look-ahead Bias（数据泄露）',
             1.0, 2.75, 5.2, 1.6, font_size=12, color=PPT_GRAY)

    add_rect(slide, 7.0, 2.1, 5.6, 2.4, fill_color=PPTColor(0x00,0x28,0x14))
    add_text(slide, '✅  新版（Walk-Forward）', 7.2, 2.2, 5, 0.5, font_size=13, bold=True, color=PPT_GREEN)
    add_text(slide, '前 60% 训练 → 预测后面\n每 10% 滚动重新训练\n每个预测点只用历史数据\n= 真实部署环境的结果',
             7.2, 2.75, 5.2, 1.6, font_size=12, color=PPT_GRAY)

    # 流程图文字版
    steps = ['历史数据 60%', '→ 训练 HMM', '→ 预测下一段', '→ 滚动前进 10%', '→ 重复']
    for i, s in enumerate(steps):
        x = 0.7 + i * 2.45
        clr = PPT_GREEN if i % 2 == 0 else PPT_CARD
        add_rect(slide, x, 5.0, 2.2, 0.7, fill_color=clr)
        add_text(slide, s, x+0.1, 5.05, 2.0, 0.6, font_size=11, bold=True,
                 color=PPT_WHITE if i % 2 == 0 else PPT_ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide, '结果：每一笔交易的预测，模型从未见过未来的数据',
             0.8, 6.1, 12, 0.6, font_size=13, color=PPT_ACCENT, bold=True)


def slide_signals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '04  14 信号投票', 0.8, 0.4, 10, 0.6, font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, '双重过滤，降低假信号', 0.8, 0.85, 10, 0.8, font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    add_text(slide, 'HMM 判定牛市  +  14 项指标满足 ≥ 9 项  →  入场',
             0.8, 1.9, 12, 0.5, font_size=14, color=PPT_ACCENT, bold=True)

    sigs = [
        ('趋势类', ['Price > EMA50', 'Price > EMA200', 'ADX > 25', 'MACD > Signal']),
        ('动量类', ['RSI < 90', 'Momentum > 1%', 'Stoch %K ↑ & < 80', 'Williams %R < -20']),
        ('量价类', ['Volume > SMA20', 'OBV > OBV EMA', 'CCI > 0', 'Price > BB 中轨']),
        ('风险类', ['Volatility < 6%', '距高点回撤 < 30%']),
    ]
    colors = [PPTColor(0x00,0x28,0x3C), PPTColor(0x0A,0x1E,0x3C),
              PPTColor(0x1A,0x0A,0x3C), PPTColor(0x20,0x18,0x08)]
    x_pos = [0.8, 3.9, 7.0, 10.1]
    for i, (cat, items) in enumerate(sigs):
        x = x_pos[i]
        w = 2.8 if i < 3 else 2.8
        add_rect(slide, x, 2.55, w, 4.4, fill_color=colors[i])
        add_text(slide, cat, x+0.15, 2.65, w-0.3, 0.45, font_size=12, bold=True, color=PPT_ACCENT)
        for j, sig in enumerate(items):
            add_rect(slide, x+0.15, 3.2+j*0.78, w-0.3, 0.6, fill_color=PPTColor(0x0A,0x0E,0x1A))
            add_text(slide, '● ' + sig, x+0.25, 3.25+j*0.78, w-0.5, 0.5,
                     font_size=10.5, color=PPT_WHITE)

    add_text(slide, '门槛：9 / 14 → 40% 仓位　　12 / 14 → 70%　　14 / 14 → 100%',
             0.8, 7.0, 12, 0.4, font_size=11, color=PPT_GRAY)


def slide_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '05  风险管理', 0.8, 0.4, 10, 0.6, font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, '三重保护机制', 0.8, 0.85, 10, 0.8, font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    rules = [
        ('🛑', '固定止损 -8%',  '任何一笔仓位亏损超过 8%\n立即出场，不等 HMM 反应\n防止单笔大额亏损'),
        ('📉', '状态止损',      'HMM 检测到 Bear/Crash 状态\n立即平仓，不设任何延迟\n跟随市场节奏'),
        ('⏳', '冷静期',        '出场后强制等待 48 小时\n(日线资产 2 个交易日)\n避免情绪化反复进出'),
        ('📐', '仓位管理',      '信号强度决定仓位大小\n9/14 → 40%\n14/14 → 100%'),
    ]
    for i, (icon, title, desc) in enumerate(rules):
        x = 0.8 + (i % 2) * 6.2
        y = 2.1 + (i // 2) * 2.5
        add_rect(slide, x, y, 5.8, 2.1, fill_color=PPT_CARD)
        add_rect(slide, x, y, 0.12, 2.1, fill_color=PPT_GREEN)
        add_text(slide, icon,  x+0.3, y+0.15, 0.8, 0.7, font_size=22)
        add_text(slide, title, x+1.1, y+0.15, 4.4, 0.5, font_size=14, bold=True, color=PPT_WHITE)
        add_text(slide, desc,  x+1.1, y+0.68, 4.4, 1.2, font_size=11, color=PPT_GRAY)

    add_text(slide, '杠杆 2.5×  ·  最长持仓 30天(1h) / 60日(日线)  ·  以上保护机制同时生效',
             0.8, 6.95, 12, 0.45, font_size=11, color=PPT_GRAY, align=PP_ALIGN.CENTER)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '06  回测结果', 0.8, 0.4, 10, 0.6, font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, '数据截至 2026-04-16', 0.8, 0.85, 10, 0.8, font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    assets = [
        ('🍎 AAPL', '+11.9%', '1.04', '-12.1%', '38', PPTColor(0x00,0x40,0x28)),
        ('🥇 Gold', '+80.6%', '1.51', '-21.5%', '25', PPTColor(0x3D,0x2E,0x00)),
        ('🥈 Silver','+80.7%','0.64', '-46.9%', '27', PPTColor(0x1A,0x1A,0x2E)),
    ]
    for i, (name, ret, sharpe, dd, trades, bg) in enumerate(assets):
        x = 0.8 + i * 4.2
        add_rect(slide, x, 2.0, 3.9, 4.5, fill_color=bg)
        add_text(slide, name,   x+0.2, 2.15, 3.5, 0.6, font_size=16, bold=True, color=PPT_WHITE)

        metrics = [('总收益', ret), ('夏普比率', sharpe), ('最大回撤', dd), ('交易笔数', trades)]
        for j, (label, val) in enumerate(metrics):
            y = 2.9 + j * 0.85
            add_text(slide, label, x+0.2, y,      2,   0.35, font_size=10, color=PPT_GRAY)
            c = PPT_GREEN if ('+' in val or (val.replace('.','').isdigit() and float(val) > 0)) else PPTColor(0xFF,0x52,0x52)
            if label == '交易笔数': c = PPT_WHITE
            if label == '夏普比率': c = PPT_GREEN if float(val) >= 1 else PPTColor(0xFF,0xD7,0x40)
            add_text(slide, val,   x+0.2, y+0.32, 3.5, 0.45, font_size=20, bold=True, color=c)

    add_text(slide, 'Gold 综合最优 · AAPL 最稳健 · Silver 收益高但回撤风险大',
             0.8, 6.7, 12, 0.5, font_size=13, color=PPT_ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    add_text(slide, '07  下一步', 0.8, 0.4, 10, 0.6, font_size=11, color=PPT_GREEN, bold=True)
    add_text(slide, 'Phase 2 计划', 0.8, 0.85, 10, 0.8, font_size=32, bold=True, color=PPT_WHITE)
    add_divider(slide, 1.75)

    plans = [
        ('01', '指标重要性分析',
         '量化 14 个信号各自的预测贡献度\n保留最有效的，去掉冗余的\n让策略更精简、更可解释'),
        ('02', '跨资产相关性',
         'Gold 与 Silver 高度相关\n研究是否应该错峰配置\n避免重复承担相同风险'),
        ('03', '与 SPY 对比',
         '验证策略是否只在牛市有效\n还是具备真实的 Alpha 来源\n提升策略可信度'),
        ('04', '参数优化',
         '对止损线 / 持仓上限 / 信号门槛\n做网格搜索，找到最优组合\n减少人为经验干预'),
    ]
    for i, (num, title, desc) in enumerate(plans):
        x = 0.8 + (i % 2) * 6.2
        y = 2.1 + (i // 2) * 2.5
        add_rect(slide, x, y, 5.8, 2.1, fill_color=PPT_CARD)
        add_text(slide, num,   x+0.2, y+0.15, 0.8, 0.7, font_size=26, bold=True, color=PPT_GREEN)
        add_text(slide, title, x+1.1, y+0.15, 4.4, 0.5, font_size=14, bold=True, color=PPT_WHITE)
        add_text(slide, desc,  x+0.2, y+0.72, 5.3, 1.2, font_size=10.5, color=PPT_GRAY)


def slide_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PPT_BG)
    add_rect(slide, 0, 0, 0.06, 7.5, fill_color=PPT_GREEN)

    # 装饰
    add_rect(slide, 6.5, 0, 6.83, 7.5, fill_color=PPT_CARD)
    circle = slide.shapes.add_shape(9, Inches(7.5), Inches(2.0), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PPTColor(0x00, 0x28, 0x14)
    circle.line.fill.background()

    add_text(slide, 'Thank You', 0.8, 2.2, 5.5, 1.2, font_size=48, bold=True, color=PPT_WHITE)
    add_divider(slide, 3.55)
    add_text(slide, '仅供学习研究\n不构成任何投资建议', 0.8, 3.75, 5.5, 1.2,
             font_size=16, color=PPT_GRAY, italic=True)

    add_text(slide, '📊 Dashboard', 7.2, 2.5, 5.5, 0.6, font_size=14, bold=True, color=PPT_ACCENT)
    add_text(slide, 'share.streamlit.io', 7.2, 3.1, 5.5, 0.5, font_size=12, color=PPT_GRAY)
    add_text(slide, '📁 GitHub', 7.2, 3.9, 5.5, 0.6, font_size=14, bold=True, color=PPT_ACCENT)
    add_text(slide, 'github.com/Anselzwx/hmm-trading-strategy', 7.2, 4.5, 5.5, 0.5,
             font_size=12, color=PPT_GRAY)


def build_ppt():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_hmm(prs)
    slide_walkforward(prs)
    slide_signals(prs)
    slide_risk(prs)
    slide_results(prs)
    slide_next(prs)
    slide_end(prs)

    out = os.path.join(OUT_DIR, 'HMM_Strategy_Deck.pptx')
    prs.save(out)
    print(f'✅ PPT 已保存：{out}')


# ══════════════════════════════════════════════════════════════
if __name__ == '__main__':
    build_word()
    build_ppt()
    print('\n🎉 全部完成！')
